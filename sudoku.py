from pptx import Presentation
from pptx.util import Inches
from pptx.util import Cm
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from tkinter import *
from tkinter import messagebox
from tkinter import filedialog
from pptx.oxml.xmlchemy import OxmlElement
from pptx_tools import utils


def SubElement(parent, tagname, **kwargs):
    element = OxmlElement(tagname)
    element.attrib.update(kwargs)
    parent.append(element)
    return element

def _set_top_cell_border(cell, border_color="000000", border_width='90800'):
    
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()

    lnT = SubElement(
        tcPr, 'a:lnT', w=border_width, cap='flat', cmpd='sng', algn='ctr')
    solidFill = SubElement(lnT, 'a:solidFill')
    srgbClr = SubElement(solidFill, 'a:srgbClr', val=border_color)
    return cell

def _set_bottom_cell_border(cell, border_color="000000", border_width='90800'):
    
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()

    lnB = SubElement(
        tcPr, 'a:lnB', w=border_width, cap='flat', cmpd='sng', algn='ctr')
    solidFill = SubElement(lnB, 'a:solidFill')
    srgbClr = SubElement(solidFill, 'a:srgbClr', val=border_color)
    return cell

def _set_left_cell_border(cell, border_color="000000", border_width='90800'):
    
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()

    lnL = SubElement(
        tcPr, 'a:lnL', w=border_width, cap='flat', cmpd='sng', algn='ctr')
    solidFill = SubElement(lnL, 'a:solidFill')
    srgbClr = SubElement(solidFill, 'a:srgbClr', val=border_color)
    return cell

def _set_right_cell_border(cell, border_color="000000", border_width='90800'):
    
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()

    lnR = SubElement(
        tcPr, 'a:lnR', w=border_width, cap='flat', cmpd='sng', algn='ctr')
    solidFill = SubElement(lnR, 'a:solidFill')
    srgbClr = SubElement(solidFill, 'a:srgbClr', val=border_color)
    return cell

def move_slide(old_index, new_index):
        xml_slides = presentation.slides._sldIdLst  # pylint: disable=W0212
        slides = list(xml_slides)
        xml_slides.remove(slides[old_index])
        xml_slides.insert(new_index, slides[old_index])
     
root = Tk()
root.withdraw()
sudokulocation = filedialog.askopenfilename(filetypes=[("Text files", "*.txt")])

 
presentation = Presentation()
presentation.slide_width = Inches(12)
presentation.slide_height = Inches(12)

fileObj = [line for line in open(sudokulocation, "r").read().splitlines() if line]
words = fileObj

array_length = len(words)
for z in range(array_length):

    currentpuzzle = words[z]
    layout = presentation.slide_masters[0].slide_layouts[6]
    slide = presentation.slides.add_slide(layout)
    Top = (0 + Inches(1))
    Left = (0 + Inches(1))
    Width = (presentation.slide_width - (Left * 2))
    Height = (presentation.slide_height - (Top * 2))
    SudokuW = 9
    SudokuH = 9
               

    if (z % 2) == 0:
        title = "Sudoku "
        number = str(int((z/2) + 1))
    else:
        title = "Sudoku Solution "
        number = str(int((z/2) + 0.5))
    if z==0:
        title = "Sudoku "
        number = str(z+1)
    if z==1:
        title = "Sudoku Solution "
        number = str(z)

    txBox = slide.shapes.add_textbox((0 + Inches(1)), (0 + Inches(0.1)), (presentation.slide_width - (Left * 2)), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title+number
    p.font.size = Pt(44)
    p.alignment = PP_ALIGN.CENTER


    
    shape = slide.shapes.add_table(SudokuW, SudokuH, Top, Left, Width, Height)

    tbl =  shape._element.graphic.graphicData.tbl
    style_id = '{5940675A-B579-460E-94D1-54222C63F5DA}'

    tbl[0][-1].text = style_id

    c = 0
    for i in range(SudokuW):
        for j in range(SudokuH):
            if currentpuzzle[c] == ".":
                shape.table.cell(i,j).text == " "
                if j==0:
                    cell = shape.table.cell(i,j)
                    cell = _set_left_cell_border(cell)
                elif j==2:
                    cell = shape.table.cell(i,j)
                    cell = _set_right_cell_border(cell)
                elif j==5:
                    cell = shape.table.cell(i,j)
                    cell = _set_right_cell_border(cell)
                elif j==8:
                    cell = shape.table.cell(i,j)
                    cell = _set_right_cell_border(cell)
                if i ==0:
                    cell = shape.table.cell(i,j)
                    cell = _set_top_cell_border(cell)
                elif i==2:
                    cell = shape.table.cell(i,j)
                    cell = _set_bottom_cell_border(cell)
                elif i==5:
                    cell = shape.table.cell(i,j)
                    cell = _set_bottom_cell_border(cell)
                elif i==8:
                    cell = shape.table.cell(i,j)
                    cell = _set_bottom_cell_border(cell)
                
                
                    
            else:
                shape.table.cell(i,j).text = currentpuzzle[c]
                shape.table.cell(i,j).text_frame.paragraphs[0].font.size = Pt(54)
                shape.table.cell(i,j).text_frame.paragraphs[0].margin_bottom = 0
                shape.table.cell(i,j).text_frame.paragraphs[0].margin_top = 0
                shape.table.cell(i,j).vertical_anchor = MSO_ANCHOR.MIDDLE
                shape.table.cell(i,j).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

                if j==0:
                    cell = shape.table.cell(i,j)
                    cell = _set_left_cell_border(cell)
                elif j==2:
                    cell = shape.table.cell(i,j)
                    cell = _set_right_cell_border(cell)
                elif j==5:
                    cell = shape.table.cell(i,j)
                    cell = _set_right_cell_border(cell)
                elif j==8:
                    cell = shape.table.cell(i,j)
                    cell = _set_right_cell_border(cell)
                if i==0:
                    cell = shape.table.cell(i,j)
                    cell = _set_top_cell_border(cell)
                elif i==2:
                    cell = shape.table.cell(i,j)
                    cell = _set_bottom_cell_border(cell)
                elif i==5:
                    cell = shape.table.cell(i,j)
                    cell = _set_bottom_cell_border(cell)
                elif i==8:
                    cell = shape.table.cell(i,j)
                    cell = _set_bottom_cell_border(cell)
                
                

                
            c = c + 1
for b in reversed(range(array_length)):
    if (b % 2) != 0:
        move_slide(b, array_length)

for t in reversed(range(array_length)):
    if t >= (array_length / 2):
        move_slide(t, array_length)


filename = filedialog.asksaveasfilename(filetypes=[("PowerPoint Presentation", "*.pptx")])
filetype = ".pptx"
saveaspptx = filename+filetype
presentation.save(saveaspptx)

saveaspptx = '\\'.join(saveaspptx.split('/'))

pptxfolder = saveaspptx.rsplit('\\', 1)[0]

png_folder = pptxfolder
pptx_file = saveaspptx


utils.save_pptx_as_png(png_folder, pptx_file, overwrite_folder=True)



App = Tk() 
App.withdraw()


messagebox.showinfo('Completed!', 'Your import of Sudoku puzzles into PowerPoint has completed successfully!')

App.mainloop()
os._exit()
